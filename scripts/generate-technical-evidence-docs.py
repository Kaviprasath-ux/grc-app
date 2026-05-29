"""
Generate the Technical Evidence platform documentation set:

Outputs (in docs/platform-docs/):
  TechnicalEvidence-USP.docx                   — sales / positioning brief
  TechnicalEvidence-Technical-Documentation.docx — functional/technical guide
  TechnicalEvidence-Presentation.pptx          — 12-slide pitch deck

Run:  python scripts/generate-technical-evidence-docs.py
"""

from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt, Emu
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- Brand palette ----------

BRAND_PRIMARY = RGBColor(0x0B, 0x3D, 0x91)
BRAND_ACCENT  = RGBColor(0x1E, 0x88, 0xE5)
TEXT_DARK     = RGBColor(0x1F, 0x2A, 0x44)
TEXT_MUTED    = RGBColor(0x55, 0x5F, 0x77)
TABLE_HEADER_BG = "0B3D91"
TABLE_ALT_BG    = "F2F5FB"

PPT_BRAND_PRIMARY = PPTRGBColor(0x0B, 0x3D, 0x91)
PPT_BRAND_ACCENT  = PPTRGBColor(0x1E, 0x88, 0xE5)
PPT_BRAND_DARK    = PPTRGBColor(0x0A, 0x18, 0x35)
PPT_TEXT_DARK     = PPTRGBColor(0x1F, 0x2A, 0x44)
PPT_TEXT_MUTED    = PPTRGBColor(0x55, 0x5F, 0x77)
PPT_TEXT_LIGHT    = PPTRGBColor(0xF4, 0xF7, 0xFC)
PPT_WHITE         = PPTRGBColor(0xFF, 0xFF, 0xFF)
PPT_ACCENT_GOLD   = PPTRGBColor(0xFF, 0xB7, 0x00)

SLIDE_W = PPTInches(13.333)
SLIDE_H = PPTInches(7.5)


# ============================================================
# DOCX helpers
# ============================================================

def set_cell_bg(cell, hex_color: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tc_pr.append(shd)


def add_horizontal_rule(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "8")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "1E88E5")
    p_bdr.append(bottom)
    p_pr.append(p_bdr)


def configure_styles(doc: Document):
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.font.color.rgb = TEXT_DARK
    for h_name, size, color in [
        ("Heading 1", 22, BRAND_PRIMARY),
        ("Heading 2", 16, BRAND_PRIMARY),
        ("Heading 3", 13, BRAND_ACCENT),
    ]:
        s = styles[h_name]
        s.font.name = "Calibri"
        s.font.size = Pt(size)
        s.font.bold = True
        s.font.color.rgb = color


def add_cover(doc: Document, platform: str, doc_type: str, subtitle: str):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("\n\n\n")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(platform.upper())
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = BRAND_ACCENT

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(doc_type)
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = BRAND_PRIMARY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_horizontal_rule(p)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(subtitle)
    r.font.size = Pt(14)
    r.font.italic = True
    r.font.color.rgb = TEXT_MUTED

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("\n\nGlimmora International\ninfo@glimmora.ai")
    r.font.size = Pt(11)
    r.font.color.rgb = TEXT_MUTED

    doc.add_page_break()


def add_bullets(doc: Document, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)


def add_table(doc: Document, headers, rows, col_widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = ""
        p = hdr[i].paragraphs[0]
        r = p.add_run(h)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.size = Pt(11)
        set_cell_bg(hdr[i], TABLE_HEADER_BG)
        hdr[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    for idx, row_data in enumerate(rows):
        row = table.add_row().cells
        for i, v in enumerate(row_data):
            row[i].text = ""
            p = row[i].paragraphs[0]
            r = p.add_run(str(v))
            r.font.size = Pt(10.5)
            if idx % 2 == 1:
                set_cell_bg(row[i], TABLE_ALT_BG)
    if col_widths:
        for row in table.rows:
            for i, w in enumerate(col_widths):
                row.cells[i].width = Cm(w)
    return table


def add_callout(doc: Document, title: str, body: str):
    table = doc.add_table(rows=1, cols=1)
    cell = table.rows[0].cells[0]
    set_cell_bg(cell, TABLE_ALT_BG)
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(title + "\n")
    r.font.bold = True
    r.font.color.rgb = BRAND_PRIMARY
    r.font.size = Pt(11)
    r2 = p.add_run(body)
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = TEXT_DARK


def add_para(doc: Document, text: str, bold: bool = False, italic: bool = False):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.font.bold = bold
    r.font.italic = italic
    return p


# ============================================================
# CONTENT — USP
# ============================================================

def build_usp(doc: Document):
    add_cover(doc, "Technical Evidence Platform", "Unique Selling Proposition",
              "Continuous, audit-grade compliance evidence — automatically collected from the tools your security team already uses.")

    doc.add_heading("1. The Problem We Solve", level=1)
    add_para(doc,
        "Compliance teams burn weeks every audit cycle chasing the same artefacts: user lists "
        "from Azure Entra, vulnerability reports from Tenable and Acunetix, change tickets "
        "from ServiceNow, firewall configurations from Palo Alto and FortiGate. The data lives "
        "in nine different consoles, each with its own export format, refresh cadence, and "
        "access ritual. By the time the evidence reaches the auditor, it is stale, screenshot-"
        "based, and impossible to trust as a current control statement.")
    add_para(doc,
        "Glimmora Technical Evidence eliminates this drudgery. Connect each platform once, "
        "and the system pulls evidence on demand, hashes it, and publishes it directly into "
        "your compliance Evidence library with control mappings already suggested by AI. The "
        "auditor opens GRC, sees a CSV with a timestamp and a SHA-256 fingerprint, and the "
        "conversation is over.")

    doc.add_heading("2. What Makes It Unique", level=1)

    doc.add_heading("2.1 Nine production-grade connectors out of the box", level=2)
    add_bullets(doc, [
        "Azure Entra (identity + access reviews via Microsoft Graph)",
        "Azure Security (Defender, Sentinel, CSPM)",
        "Google Cloud Security Command Center",
        "ServiceNow (incidents, change tickets, asset inventory)",
        "Tenable (vulnerability scans, asset coverage)",
        "Acunetix (web-app vulnerability findings)",
        "Cisco ASA / FTD (firewall configuration & rule sets)",
        "Palo Alto Networks (NGFW configuration & threat logs)",
        "FortiGate (UTM configuration & policy)",
    ])
    add_callout(doc, "Why this matters",
        "Competing tools require you to upload CSVs by hand or hire a consultant to wire up a "
        "custom integration per platform. We ship the connectors as part of the product — onboarding "
        "is a 90-second credential-paste exercise, not a six-week SI engagement.")

    doc.add_heading("2.2 Evidence is auto-published to GRC, not parked in a silo", level=2)
    add_para(doc,
        "Most evidence-collection tools dump CSVs into their own object store and leave the user "
        "to wire it into the compliance program manually. Technical Evidence is built into the "
        "same database as the GRC Evidence library: every successful collection creates or updates "
        "an Evidence record with the original CSV attached, the data-source metadata embedded, "
        "and AI-suggested links to the controls it most likely satisfies. One click confirms the "
        "mapping; the auditor sees a fully-traceable evidence chain from regulator-named control "
        "down to the row in the connector's response.")

    doc.add_heading("2.3 Credentials are vaulted, never seen by humans", level=2)
    add_bullets(doc, [
        "All secrets stored AES-256-GCM encrypted (12-byte IV + 16-byte auth tag) under the customer-scoped CREDENTIAL_ENCRYPTION_KEY.",
        "Encryption is performed in the API layer on write; the database row stores only ciphertext.",
        "Credential test calls run on demand only; results are recorded with timestamp and last-status badge in the Credential Vault.",
        "Field-level RBAC blocks credential view/edit to CustomerAdministrator or GRCAdministrator only — auditors and end users never touch them.",
    ])

    doc.add_heading("2.4 Deduplication and integrity by design", level=2)
    add_para(doc,
        "Every record is hashed (SHA-256 of canonical JSON) before insert, so re-collecting the "
        "same data does not bloat storage and tamper attempts are detectable. Each collection "
        "snapshot keeps an immutable record set tied to its `collectedAt` timestamp, giving "
        "auditors a defensible 'this is what the world looked like at 14:32 UTC on the 12th' "
        "answer for any control check.")

    doc.add_heading("2.5 Pay only for what you use — independent of GRC", level=2)
    add_para(doc,
        "Technical Evidence is a standalone purchasable module. Customers running a third-party "
        "GRC product can still buy Technical Evidence and export evidence on demand. Customers "
        "running Glimmora GRC get the deepest experience because evidence flows natively into "
        "their controls library, but the platform is never a forced bundle.")

    doc.add_heading("3. Quantified Value", level=1)
    add_table(doc,
        ["Outcome", "Before Technical Evidence", "With Technical Evidence"],
        [
            ["Evidence freshness", "Quarterly screenshot drops", "On-demand, with timestamp + hash"],
            ["Time per evidence cycle", "3–6 hours per control owner", "Sub-5 minutes per data source"],
            ["Auditor questions about provenance", "10–20 per engagement", "Effectively zero — provenance is built in"],
            ["Annual evidence storage cost", "Drift toward duplicates", "SHA-256 dedup eliminates duplicates"],
            ["Onboarding effort per new tool", "1–6 weeks (custom integration)", "60–90 seconds (paste credential, test, collect)"],
        ],
        col_widths=[5, 5.5, 5.5],
    )

    doc.add_heading("4. Who Buys This", level=1)
    add_bullets(doc, [
        "CISOs and Heads of GRC who need defensible evidence for ISO 27001, SOC 2, PCI DSS, HIPAA, RBI / SAMA / NESA / NCA audits.",
        "Compliance leads tired of being the human glue between IT operations and audit.",
        "Internal audit functions that want to spot-check controls between formal engagements.",
        "MSSPs and consulting firms running compliance-as-a-service for multiple end customers, billing per integration.",
    ])

    doc.add_heading("5. Why Now", level=1)
    add_para(doc,
        "Regulatory bodies across the Gulf, India, and EU are converging on continuous-compliance "
        "expectations — controls must be operating effectively, not just designed. Static evidence "
        "is being challenged in audit by automated systems that prove a control was working on the "
        "exact day of an incident. Technical Evidence is built for that bar.")

    add_callout(doc, "Bottom line",
        "Buy once, connect nine platforms in an hour, eliminate quarterly evidence-collection "
        "marathons forever. Spend that capacity on actual risk reduction instead.")


# ============================================================
# CONTENT — Functional / Technical Documentation
# ============================================================

def build_functional_doc(doc: Document):
    add_cover(doc, "Technical Evidence Platform", "Functional Documentation",
              "End-to-end capability description: data model, workflows, integrations, security, and operations.")

    doc.add_heading("1. Platform Overview", level=1)
    add_para(doc,
        "Technical Evidence is the fourth platform in the Glimmora suite, alongside GRC, TPRM, "
        "and Internal Audit. It provides a credential vault for nine third-party security and "
        "IT-operations platforms, an on-demand data-collection engine, and an evidence-publication "
        "pipeline that auto-creates compliance Evidence records inside GRC. Each customer's "
        "subscription, data, and credentials are tenant-isolated.")

    add_callout(doc, "Module gating",
        "Access is controlled by the CustomerAccount.isTechnicalEvidenceEnabled flag plus a "
        "Subscription / ModuleSubscription row with moduleCode='TECHNICAL_EVIDENCE'. The layout "
        "gate blocks deep-links to /technical-evidence/* if either is missing and redirects to "
        "/subscription-required.")

    doc.add_heading("2. Capability Map", level=1)
    add_table(doc,
        ["Capability", "Where", "What it does"],
        [
            ["Credential Vault", "/technical-evidence/settings", "Add/edit/delete encrypted credentials for the 9 supported platforms."],
            ["Per-platform configuration", "/technical-evidence/settings/[platform]", "Dynamic form populated from PLATFORM_INFO with the platform's required fields and authType."],
            ["Connection test", "/api/integration-credentials/[id]/test", "On-demand reachability check; updates lastTestedAt / lastTestStatus on the credential row."],
            ["Data-source bootstrap", "/api/technical-evidence/bootstrap", "After a credential is saved, creates one TechnicalEvidenceCollection per data source exposed by the connector."],
            ["Dashboard", "/technical-evidence/dashboard", "Cards per collection with status, record count, last-collected timestamp, and Collect/View buttons."],
            ["Data collection", "/api/technical-evidence/[id]/collect", "Decrypts credential, calls connector.collectData(), inserts records (batched, hashed), auto-publishes Evidence."],
            ["Control mapping", "/api/technical-evidence/[id]/control-mappings", "AI-suggested and user-confirmed mapping between a collection and one or more controls/frameworks."],
            ["Evidence publication", "/api/technical-evidence/[id]/publish", "Creates / updates an Evidence row with the collection's CSV attached and control links populated."],
            ["Account overview (admin)", "/technical-evidence/account-overview", "Superadmin-only customer list with TE enabled, primary admin, credential and collection counts."],
        ],
        col_widths=[4.5, 5, 6.5],
    )

    doc.add_heading("3. Data Model", level=1)

    doc.add_heading("3.1 IntegrationCredential", level=2)
    add_para(doc,
        "Encrypted credentials for a third-party platform. One row per platform + named instance "
        "per tenant (unique constraint on customerAccountId + platform + name). Stores ciphertext "
        "only — the encryption library is src/lib/credential-encryption.ts.")

    doc.add_heading("3.2 TechnicalEvidenceCollection", level=2)
    add_para(doc,
        "A definition of one data source backed by a specific credential. For example, the Azure "
        "Entra credential might back collections for 'All Users', 'Dormant Users', 'Privileged "
        "Roles', and 'Conditional Access Policies'. Lifecycle status field tracks the last collection "
        "attempt (never / collecting / success / failed) and a record count.")

    doc.add_heading("3.3 TechnicalEvidenceRecord", level=2)
    add_para(doc,
        "Individual rows returned from a connector, stored as JSON. Each record has a SHA-256 "
        "hash so re-collection is idempotent (same data → same hash → upsert is a no-op).")

    doc.add_heading("3.4 TechnicalEvidenceControlMapping", level=2)
    add_para(doc,
        "Many-to-one link from a collection to a compliance Control, with an isAISuggested / "
        "isUserConfirmed flag and an optional confidence score. The AI suggestion engine uses "
        "keyword and semantic match against the customer's existing control library; the user "
        "confirms or rejects in one click.")

    doc.add_heading("3.5 Evidence cross-reference", level=2)
    add_para(doc,
        "The GRC Evidence model has a nullable technicalEvidenceCollectionId field. When set, the "
        "Evidence row was auto-created by the publish pipeline and carries a back-reference to "
        "the collection that produced it. Re-collection updates the same Evidence row in place "
        "rather than creating duplicates.")

    doc.add_heading("4. End-to-End Workflow", level=1)

    add_table(doc,
        ["Step", "Actor", "Action", "System effect"],
        [
            ["1", "CustomerAdmin", "Opens /technical-evidence/settings", "Renders 9 platform cards from PLATFORM_INFO."],
            ["2", "CustomerAdmin", "Clicks a card → fills credential form", "POST /api/integration-credentials. Encrypts with AES-256-GCM. Inserts IntegrationCredential row."],
            ["3", "System", "Bootstrap fires", "POST /api/technical-evidence/bootstrap. For each data source exposed by the connector, upserts a TechnicalEvidenceCollection."],
            ["4", "CustomerAdmin", "Clicks 'Test connection'", "POST /api/integration-credentials/[id]/test. Decrypts in memory, calls connector.testConnection(). Updates lastTestStatus."],
            ["5", "CustomerAdmin", "Opens dashboard, clicks 'Collect' on a card", "POST /api/technical-evidence/[id]/collect. Updates status='collecting'. Decrypts credential. Calls connector.collectData()."],
            ["6", "System", "Persists records", "Deletes prior TechnicalEvidenceRecord rows for that collection. Bulk-inserts new records in batches of 100, each hashed."],
            ["7", "System", "Suggests control mappings", "Keyword/semantic match against customer's control library. Inserts TechnicalEvidenceControlMapping rows with suggestedByAI=true."],
            ["8", "System", "Auto-publishes Evidence", "publishTechnicalEvidence(): creates or updates Evidence with the CSV attached, links to confirmed/suggested controls."],
            ["9", "CustomerAdmin / Reviewer", "Reviews suggestions in the collection detail page", "PATCH /api/technical-evidence/[id]/control-mappings confirms each mapping (or POSTs a manual one)."],
            ["10", "Auditor", "Opens GRC > Evidence", "Sees Evidence rows tagged 'Auto-collected', timestamped, with the original CSV attached."],
        ],
        col_widths=[1.2, 3, 5.5, 6.3],
    )

    doc.add_heading("5. Supported Integrations", level=1)
    add_table(doc,
        ["Platform", "Auth", "Sample data sources"],
        [
            ["Azure Entra (Active Directory)", "OAuth2 (client credentials)", "All users; dormant users; privileged roles; conditional access policies; sign-in risk events."],
            ["Azure Security", "OAuth2 (client credentials)", "Defender for Cloud recommendations; Sentinel incidents; CSPM findings; secure score history."],
            ["Google Cloud", "Service account JSON", "SCC findings; IAM bindings; org policies; asset inventory."],
            ["ServiceNow", "Basic auth", "Incidents; change tickets; CMDB asset list; security incidents."],
            ["Tenable", "API key + secret", "Vulnerability scans; asset coverage; plugin findings."],
            ["Acunetix", "API key", "Web-app vulnerabilities; scan history; target inventory."],
            ["Cisco (ASA / FTD)", "Basic auth", "Firewall rule set; object groups; AAA configuration; running-config snapshot."],
            ["Palo Alto Networks", "API key", "NGFW security rules; threat log summary; address objects."],
            ["FortiGate", "API token", "Firewall policies; UTM profiles; address book; admin accounts."],
        ],
        col_widths=[5, 4, 7],
    )

    doc.add_heading("6. Security & Compliance", level=1)

    doc.add_heading("6.1 Credential encryption", level=2)
    add_bullets(doc, [
        "Algorithm: AES-256-GCM with 12-byte random IV + 16-byte auth tag.",
        "Key: 32-byte symmetric key derived from CREDENTIAL_ENCRYPTION_KEY (env var, stored in DigitalOcean encrypted env). Separate from the field-encryption key used for Bytes columns.",
        "Storage: ciphertext only on IntegrationCredential.credentialsEncrypted. Plaintext exists only in process memory during a collection or test call.",
        "Rotation: rotating CREDENTIAL_ENCRYPTION_KEY requires a re-encrypt sweep (procedure documented in SECURITY.md).",
    ])

    doc.add_heading("6.2 RBAC", level=2)
    add_table(doc,
        ["Resource", "Required permission", "Roles with default access"],
        [
            ["Credential CRUD", "technical-evidence.integration-credentials:create / edit / delete", "GRCAdministrator, CustomerAdministrator (TE module)"],
            ["Dashboard", "technical-evidence.dashboard:view", "GRCAdministrator, CustomerAdministrator (TE module)"],
            ["Collection trigger", "technical-evidence.dashboard:create", "GRCAdministrator, CustomerAdministrator (TE module)"],
            ["Mapping edit", "technical-evidence.dashboard:edit", "GRCAdministrator, CustomerAdministrator (TE module)"],
            ["Customer overview", "technical-evidence.account-overview:view", "GRCAdministrator (superadmin only)"],
            ["Organization profile", "technical-evidence.organization-profile:view / edit", "GRCAdministrator, CustomerAdministrator (TE module)"],
        ],
        col_widths=[5.5, 6.5, 5],
    )

    doc.add_heading("6.3 Multi-tenant isolation", level=2)
    add_bullets(doc, [
        "Every IntegrationCredential and TechnicalEvidenceCollection row is scoped by customerAccountId.",
        "API routes wrap handlers with withAuth(), which resolves the session, extracts the tenant filter, and rejects cross-tenant reads at the query layer.",
        "Subscription gating (Subscription + ModuleSubscription) is checked on every layout render and on every navigation that resolves to /technical-evidence/*.",
    ])

    doc.add_heading("7. Operational Notes", level=1)
    add_bullets(doc, [
        "Connectors are stateless — each collect call is fully self-contained. No background pollers; the system collects only when invoked by a user or by a scheduled cron (where configured).",
        "Records are written in batches of 100 to avoid Postgres parameter-count limits on bulk inserts.",
        "Connection-test results are advisory; a failed test still allows collection (some platforms throttle test endpoints separately from data endpoints).",
        "When a credential is deleted, all child collections cascade-delete (Prisma onDelete: Cascade). The auto-published Evidence row remains but its technicalEvidenceCollectionId becomes null.",
    ])

    doc.add_heading("8. Roadmap (next 2 quarters)", level=1)
    add_bullets(doc, [
        "Scheduled collection (cron per collection, defined by the customer admin).",
        "Diff-based collection (only push changed records into the Evidence CSV).",
        "Additional connectors: CrowdStrike Falcon, Wiz, Okta, AWS Security Hub.",
        "Bring-your-own-LLM for control-mapping suggestions (currently a single fixed model).",
        "Evidence-level approval workflow (a reviewer must confirm before the auto-published Evidence is visible to auditors).",
    ])

    doc.add_heading("9. Glossary", level=1)
    add_table(doc,
        ["Term", "Meaning"],
        [
            ["Connector", "TypeScript module under src/lib/integrations/ that implements testConnection, getAvailableDataSources, and collectData for one platform."],
            ["Collection", "A definition of one data source backed by one credential — e.g. 'Azure Entra: Dormant Users'."],
            ["Data source", "A specific export from a connector — e.g. all users vs. just dormant users."],
            ["Record", "A single row collected from a data source, stored as JSON with a SHA-256 hash."],
            ["Mapping", "A link between a collection and a control, marking that the collection serves as evidence for that control."],
        ],
        col_widths=[4, 13],
    )


# ============================================================
# PPTX helpers
# ============================================================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, italic=False, color=PPT_TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = PPTPt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets_ppt(slide, left, top, width, height, items, *,
                    size=18, color=PPT_TEXT_DARK, bullet_color=None, line_spacing=1.15):
    bc = bullet_color or PPT_BRAND_ACCENT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = PPTPt(6)
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = "Calibri"
        r1.font.size = PPTPt(size)
        r1.font.bold = True
        r1.font.color.rgb = bc
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = PPTPt(size)
        r2.font.color.rgb = color
    return tb


def add_header(slide, title, eyebrow=None):
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, PPTInches(0.18), PPT_BRAND_ACCENT)
    if eyebrow:
        add_text(slide, PPTInches(0.6), PPTInches(0.35), PPTInches(10), PPTInches(0.35),
                 eyebrow.upper(), size=11, bold=True, color=PPT_BRAND_ACCENT)
    add_text(slide, PPTInches(0.6), PPTInches(0.7), PPTInches(12.1), PPTInches(0.8),
             title, size=30, bold=True, color=PPT_BRAND_PRIMARY)
    add_rect(slide, PPTInches(0.6), PPTInches(1.55), PPTInches(1.5), PPTInches(0.05), PPT_ACCENT_GOLD)


def add_footer(slide, page_num, total):
    add_rect(slide, Emu(0), SLIDE_H - PPTInches(0.35), SLIDE_W, PPTInches(0.35), PPT_BRAND_DARK)
    add_text(slide, PPTInches(0.5), SLIDE_H - PPTInches(0.32), PPTInches(8), PPTInches(0.3),
             "Glimmora International  •  Technical Evidence  •  info@glimmora.ai",
             size=10, color=PPT_TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - PPTInches(1.2), SLIDE_H - PPTInches(0.32), PPTInches(0.9), PPTInches(0.3),
             f"{page_num} / {total}",
             size=10, color=PPT_TEXT_LIGHT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# CONTENT — Presentation
# ============================================================

def build_presentation(prs: Presentation):
    blank = prs.slide_layouts[6]

    # ----- Slide 1: Cover -----
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PPT_BRAND_DARK)
    add_rect(slide, Emu(0), Emu(0), PPTInches(0.6), SLIDE_H, PPT_BRAND_ACCENT)
    add_rect(slide, PPTInches(0.6), Emu(0), SLIDE_W - PPTInches(0.6), PPTInches(0.2), PPT_ACCENT_GOLD)
    add_text(slide, PPTInches(1.2), PPTInches(2.2), PPTInches(11), PPTInches(0.6),
             "GLIMMORA TECHNICAL EVIDENCE",
             size=14, bold=True, color=PPT_BRAND_ACCENT)
    add_text(slide, PPTInches(1.2), PPTInches(2.8), PPTInches(11), PPTInches(1.8),
             "Continuous, Audit-Grade Compliance Evidence",
             size=44, bold=True, color=PPT_WHITE)
    add_text(slide, PPTInches(1.2), PPTInches(4.6), PPTInches(11), PPTInches(1.3),
             "Automatically collected from the tools your security team already uses.",
             size=18, color=PPT_TEXT_LIGHT)
    add_text(slide, PPTInches(1.2), PPTInches(6.6), PPTInches(11), PPTInches(0.4),
             "Glimmora International  •  info@glimmora.ai",
             size=12, color=PPT_TEXT_LIGHT)

    TOTAL = 12

    # ----- Slide 2: The Problem -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Evidence collection is the silent killer of audit cycles", "the problem")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "Compliance teams spend 3–6 hours per control owner per audit, every audit, chasing the same artefacts.",
        "Evidence lives across nine consoles — Azure, GCP, ServiceNow, Tenable, Acunetix, Cisco, Palo Alto, FortiGate, custom.",
        "Screenshots and CSV exports are stale the moment they're attached; auditors challenge provenance.",
        "Every new tool the security team adopts is another six-week custom integration project for GRC.",
        "Regulators are converging on continuous-compliance expectations: 'designed' is no longer enough; controls must be operating.",
    ], size=18)
    add_footer(slide, 2, TOTAL)

    # ----- Slide 3: The Solution -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Connect once. Collect on demand. Publish into GRC automatically.", "the solution")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "A vault for nine pre-built connectors — 90-second paste-and-test setup per platform.",
        "On-demand data collection: click a card, the connector fetches live data, records are hashed and stored.",
        "AI suggests which control(s) the collection maps to. One click confirms.",
        "Evidence is auto-published into the GRC Evidence library with the original CSV attached and a full provenance trail.",
        "Credentials are AES-256-GCM encrypted; raw secrets never appear in logs or UI.",
    ], size=18)
    add_footer(slide, 3, TOTAL)

    # ----- Slide 4: Architecture -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "How it works", "architecture")
    flow_top = PPTInches(2.4)
    step_w = PPTInches(2.3)
    step_h = PPTInches(2.1)
    gap = PPTInches(0.2)
    start_left = PPTInches(0.6)
    steps = [
        ("1. Vault", "Encrypted credential stored under tenant"),
        ("2. Bootstrap", "Connector exposes data sources as Collections"),
        ("3. Collect", "Decrypt → fetch → hash → batch insert"),
        ("4. Map", "AI suggests Control links"),
        ("5. Publish", "Evidence row + CSV attached in GRC"),
    ]
    cur_left = start_left
    for title, body in steps:
        add_rect(slide, cur_left, flow_top, step_w, step_h, PPT_BRAND_PRIMARY)
        add_text(slide, cur_left, flow_top + PPTInches(0.25), step_w, PPTInches(0.6), title,
                 size=18, bold=True, color=PPT_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, cur_left + PPTInches(0.1), flow_top + PPTInches(0.9), step_w - PPTInches(0.2), PPTInches(1.1), body,
                 size=12, color=PPT_TEXT_LIGHT, align=PP_ALIGN.CENTER)
        cur_left += step_w + gap
    add_text(slide, PPTInches(0.6), PPTInches(5.0), PPTInches(12), PPTInches(0.8),
             "Each step is tenant-scoped, idempotent, and auditable end-to-end.",
             size=14, italic=True, color=PPT_TEXT_MUTED)
    add_footer(slide, 4, TOTAL)

    # ----- Slide 5: Connectors -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Nine production connectors out of the box", "integrations")
    grid = [
        ("Azure Entra", "Identity, MFA, dormant users, conditional access"),
        ("Azure Security", "Defender, Sentinel, CSPM, secure score"),
        ("Google Cloud", "SCC findings, IAM, org policies, asset inventory"),
        ("ServiceNow", "Incidents, change tickets, CMDB"),
        ("Tenable", "Vulnerability scans, asset coverage"),
        ("Acunetix", "Web-app vulnerabilities, scan history"),
        ("Cisco ASA / FTD", "Firewall rules, AAA, running-config"),
        ("Palo Alto Networks", "NGFW rules, threat logs, address objects"),
        ("FortiGate", "Policies, UTM, address book, admins"),
    ]
    cols = 3
    cell_w = PPTInches(4.0)
    cell_h = PPTInches(1.2)
    grid_top = PPTInches(2.0)
    grid_left = PPTInches(0.6)
    for i, (name, desc) in enumerate(grid):
        col = i % cols
        row = i // cols
        x = grid_left + col * (cell_w + PPTInches(0.15))
        y = grid_top + row * (cell_h + PPTInches(0.2))
        add_rect(slide, x, y, cell_w, cell_h, PPTRGBColor(0xF2, 0xF5, 0xFB))
        add_text(slide, x + PPTInches(0.2), y + PPTInches(0.1), cell_w - PPTInches(0.4), PPTInches(0.45),
                 name, size=15, bold=True, color=PPT_BRAND_PRIMARY)
        add_text(slide, x + PPTInches(0.2), y + PPTInches(0.55), cell_w - PPTInches(0.4), PPTInches(0.6),
                 desc, size=11, color=PPT_TEXT_MUTED)
    add_footer(slide, 5, TOTAL)

    # ----- Slide 6: User experience -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "From credential to audit-grade Evidence in 90 seconds", "user experience")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "Open the Credential Vault. Click a platform card. Paste the credential. Test. Save.",
        "The dashboard immediately shows the collections that platform supports.",
        "Click Collect. Watch records come in, hashed and deduplicated.",
        "Open the collection detail. AI shows you which controls it likely satisfies. Confirm in one click.",
        "Switch to GRC > Evidence. The CSV is already there, mapped, timestamped, defensible.",
    ], size=18)
    add_footer(slide, 6, TOTAL)

    # ----- Slide 7: Security -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Security model", "trust")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "Credentials encrypted with AES-256-GCM under a customer-scoped CREDENTIAL_ENCRYPTION_KEY.",
        "Decryption happens only in process memory during a collect or test call — never written to disk or logs.",
        "Record hashing (SHA-256 of canonical JSON) makes re-collection idempotent and tamper-detectable.",
        "Multi-tenant isolation enforced at every query: customerAccountId is part of every read and write.",
        "RBAC restricts credential CRUD to customer-administrators; auditors and end users never see secrets.",
        "TLS-only network egress to all connector endpoints.",
    ], size=18)
    add_footer(slide, 7, TOTAL)

    # ----- Slide 8: Compliance value -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Quantified before-and-after", "value")

    headers = ["Outcome", "Before", "With Technical Evidence"]
    data = [
        ("Evidence freshness", "Quarterly screenshot drops", "On-demand, hashed, timestamped"),
        ("Hours per audit cycle", "3–6 hrs per control owner", "<5 min per data source"),
        ("Provenance questions", "10–20 per engagement", "Effectively zero"),
        ("New-tool onboarding", "1–6 weeks", "60–90 seconds"),
        ("Duplicate storage", "Drifts upward", "SHA-256 dedup"),
    ]
    table_top = PPTInches(2.0)
    row_h = PPTInches(0.55)
    col_w = [PPTInches(3.6), PPTInches(4.5), PPTInches(4.5)]
    col_x = [PPTInches(0.6), PPTInches(4.2), PPTInches(8.7)]
    for i, h in enumerate(headers):
        add_rect(slide, col_x[i], table_top, col_w[i], row_h, PPT_BRAND_PRIMARY)
        add_text(slide, col_x[i] + PPTInches(0.15), table_top + PPTInches(0.1),
                 col_w[i] - PPTInches(0.3), row_h - PPTInches(0.1),
                 h, size=14, bold=True, color=PPT_WHITE)
    for r_idx, row in enumerate(data):
        y = table_top + row_h * (r_idx + 1)
        bg = PPTRGBColor(0xF2, 0xF5, 0xFB) if r_idx % 2 == 0 else PPT_WHITE
        for i, val in enumerate(row):
            add_rect(slide, col_x[i], y, col_w[i], row_h, bg)
            add_text(slide, col_x[i] + PPTInches(0.15), y + PPTInches(0.1),
                     col_w[i] - PPTInches(0.3), row_h - PPTInches(0.1),
                     val, size=12, color=PPT_TEXT_DARK)
    add_footer(slide, 8, TOTAL)

    # ----- Slide 9: Buyer personas -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Who buys Technical Evidence", "go-to-market")
    cards = [
        ("CISO / Head of GRC", "Defensible evidence for ISO 27001, SOC 2, PCI DSS, RBI, SAMA, NESA, NCA, HIPAA."),
        ("Compliance Lead", "Stop being the human glue between IT operations and audit."),
        ("Internal Audit", "Spot-check controls between formal engagements."),
        ("MSSP / vCISO", "Compliance-as-a-service across many end customers, billable per integration."),
    ]
    card_w = PPTInches(6.0)
    card_h = PPTInches(2.0)
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = PPTInches(0.6) + col * (card_w + PPTInches(0.3))
        y = PPTInches(2.0) + row * (card_h + PPTInches(0.3))
        add_rect(slide, x, y, card_w, card_h, PPTRGBColor(0xF2, 0xF5, 0xFB))
        add_rect(slide, x, y, PPTInches(0.08), card_h, PPT_BRAND_ACCENT)
        add_text(slide, x + PPTInches(0.3), y + PPTInches(0.15), card_w - PPTInches(0.4), PPTInches(0.5),
                 title, size=18, bold=True, color=PPT_BRAND_PRIMARY)
        add_text(slide, x + PPTInches(0.3), y + PPTInches(0.7), card_w - PPTInches(0.4), PPTInches(1.2),
                 body, size=13, color=PPT_TEXT_DARK)
    add_footer(slide, 9, TOTAL)

    # ----- Slide 10: Pricing model -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Pricing & packaging", "commercials")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "Sold as an independent module — buy alone or bundled with GRC, TPRM, or Internal Audit.",
        "Per-customer flat subscription tier with a fair-use cap on collection runs per month.",
        "All nine connectors included — no per-integration upsell.",
        "Custom connector development available as a fixed-scope add-on for non-listed platforms.",
        "Grandfathering: existing Glimmora GRC customers can enable Technical Evidence with a flag flip; no migration required.",
    ], size=18)
    add_footer(slide, 10, TOTAL)

    # ----- Slide 11: Roadmap -----
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Roadmap (next two quarters)", "what's next")
    add_bullets_ppt(slide, PPTInches(0.6), PPTInches(2), PPTInches(12), PPTInches(4.5), [
        "Scheduled collection — cron per collection, defined by the customer admin.",
        "Diff-based collection — push only changed records into the published Evidence CSV.",
        "More connectors: CrowdStrike Falcon, Wiz, Okta, AWS Security Hub.",
        "Bring-your-own-LLM for control-mapping suggestions.",
        "Evidence-level approval workflow — a reviewer signs off before auto-published Evidence is auditor-visible.",
    ], size=18)
    add_footer(slide, 11, TOTAL)

    # ----- Slide 12: Close -----
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PPT_BRAND_DARK)
    add_rect(slide, Emu(0), Emu(0), PPTInches(0.6), SLIDE_H, PPT_BRAND_ACCENT)
    add_text(slide, PPTInches(1.2), PPTInches(2.5), PPTInches(11), PPTInches(1.2),
             "Audit evidence shouldn't be a quarterly fire drill.",
             size=40, bold=True, color=PPT_WHITE)
    add_text(slide, PPTInches(1.2), PPTInches(4.0), PPTInches(11), PPTInches(0.7),
             "Connect your stack once. Stay audit-ready continuously.",
             size=22, color=PPT_TEXT_LIGHT)
    add_text(slide, PPTInches(1.2), PPTInches(6.5), PPTInches(11), PPTInches(0.5),
             "info@glimmora.ai  •  glimmora.ai",
             size=14, color=PPT_TEXT_LIGHT)


# ============================================================
# Main
# ============================================================

def main():
    out_dir = Path(__file__).resolve().parents[1] / "docs" / "platform-docs" / "TechnicalEvidence"
    out_dir.mkdir(parents=True, exist_ok=True)

    # USP
    doc_usp = Document()
    configure_styles(doc_usp)
    build_usp(doc_usp)
    usp_path = out_dir / "TechnicalEvidence-USP.docx"
    doc_usp.save(str(usp_path))
    print(f"Wrote {usp_path.relative_to(Path(__file__).resolve().parents[1])}")

    # Functional doc
    doc_fn = Document()
    configure_styles(doc_fn)
    build_functional_doc(doc_fn)
    fn_path = out_dir / "TechnicalEvidence-Functional-Documentation.docx"
    doc_fn.save(str(fn_path))
    print(f"Wrote {fn_path.relative_to(Path(__file__).resolve().parents[1])}")

    # Presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_presentation(prs)
    ppt_path = out_dir / "TechnicalEvidence-Presentation.pptx"
    prs.save(str(ppt_path))
    print(f"Wrote {ppt_path.relative_to(Path(__file__).resolve().parents[1])}")


if __name__ == "__main__":
    main()

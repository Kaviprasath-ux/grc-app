"""
Generate PowerPoint (.pptx) decks for the Internal Audit and GRC platforms.

Produces 2 presentations in docs/platform-docs/:
  - InternalAudit-Presentation.pptx
  - GRC-Presentation.pptx

Run:  python scripts/generate-platform-presentations.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Brand palette ----------

BRAND_PRIMARY = RGBColor(0x0B, 0x3D, 0x91)   # deep blue
BRAND_ACCENT  = RGBColor(0x1E, 0x88, 0xE5)   # sky blue
BRAND_DARK    = RGBColor(0x0A, 0x18, 0x35)   # almost-black blue
TEXT_DARK     = RGBColor(0x1F, 0x2A, 0x44)
TEXT_MUTED    = RGBColor(0x55, 0x5F, 0x77)
TEXT_LIGHT    = RGBColor(0xF4, 0xF7, 0xFC)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GOLD   = RGBColor(0xFF, 0xB7, 0x00)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Helpers ----------

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=TEXT_DARK, bullet_color=None, line_spacing=1.15):
    bc = bullet_color or BRAND_ACCENT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet marker
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bc
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_footer(slide, page_num: int, total: int, platform: str):
    # Bottom accent bar
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), BRAND_DARK)
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.32), Inches(8), Inches(0.3),
             f"Glimmora International  •  {platform}  •  info@glimmora.ai",
             size=10, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.32), Inches(0.9), Inches(0.3),
             f"{page_num} / {total}",
             size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_header(slide, title: str, eyebrow: str = None):
    # Top accent bar (thin)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.18), BRAND_ACCENT)

    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(10), Inches(0.35),
                 eyebrow.upper(),
                 size=11, bold=True, color=BRAND_ACCENT)

    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(0.8),
             title,
             size=30, bold=True, color=BRAND_PRIMARY)

    # Title underline
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(1.5), Inches(0.05), ACCENT_GOLD)


# ---------- Slide builders ----------

def build_cover_slide(prs: Presentation, platform: str, subtitle: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BRAND_DARK)

    # Decorative diagonal accent
    accent_left = add_rect(slide, Emu(0), Emu(0), Inches(0.6), SLIDE_H, BRAND_ACCENT)
    accent_top_strip = add_rect(slide, Inches(0.6), Inches(0.0), SLIDE_W - Inches(0.6), Inches(0.2), ACCENT_GOLD)

    # Eyebrow
    add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(0.5),
             "GLIMMORA INTERNATIONAL",
             size=14, bold=True, color=ACCENT_GOLD, font="Calibri")

    # Main title
    add_text(slide, Inches(1.2), Inches(2.4), Inches(11), Inches(1.5),
             platform,
             size=56, bold=True, color=WHITE)

    # Title underline accent
    add_rect(slide, Inches(1.2), Inches(3.85), Inches(2.0), Inches(0.08), ACCENT_GOLD)

    # Subtitle
    add_text(slide, Inches(1.2), Inches(4.1), Inches(11), Inches(1.5),
             subtitle,
             size=22, color=TEXT_LIGHT, font="Calibri")

    # Bottom-right tag
    # Bottom-right contact + tagline
    add_text(slide, SLIDE_W - Inches(6.0), SLIDE_H - Inches(1.2), Inches(5.5), Inches(0.45),
             "Investor & Customer Briefing",
             size=14, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    add_text(slide, SLIDE_W - Inches(6.0), SLIDE_H - Inches(0.75), Inches(5.5), Inches(0.45),
             "info@glimmora.ai",
             size=13, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)

    return slide


def build_section_divider(prs: Presentation, label: str, title: str,
                          platform: str, page_num: int, total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BRAND_PRIMARY)

    # Big section number / label
    add_text(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.6),
             label.upper(),
             size=18, bold=True, color=ACCENT_GOLD)

    add_text(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.5),
             title,
             size=44, bold=True, color=WHITE)

    add_rect(slide, Inches(0.8), Inches(4.4), Inches(1.5), Inches(0.08), ACCENT_GOLD)

    add_footer(slide, page_num, total, platform)
    return slide


def build_content_slide(prs: Presentation, eyebrow: str, title: str, bullets: list,
                        platform: str, page_num: int, total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_header(slide, title, eyebrow=eyebrow)
    add_bullets(slide, Inches(0.7), Inches(2.0), Inches(12.0), Inches(5.0),
                bullets, size=18)
    add_footer(slide, page_num, total, platform)
    return slide


def build_two_col_slide(prs: Presentation, eyebrow: str, title: str,
                        left_heading: str, left_bullets: list,
                        right_heading: str, right_bullets: list,
                        platform: str, page_num: int, total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_header(slide, title, eyebrow=eyebrow)

    # Left column
    add_text(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(0.45),
             left_heading, size=16, bold=True, color=BRAND_PRIMARY)
    add_bullets(slide, Inches(0.7), Inches(2.55), Inches(5.8), Inches(4.4),
                left_bullets, size=15)

    # Divider
    add_rect(slide, Inches(6.65), Inches(2.0), Inches(0.03), Inches(5.0), BRAND_ACCENT)

    # Right column
    add_text(slide, Inches(6.9), Inches(2.0), Inches(5.8), Inches(0.45),
             right_heading, size=16, bold=True, color=BRAND_PRIMARY)
    add_bullets(slide, Inches(6.9), Inches(2.55), Inches(5.8), Inches(4.4),
                right_bullets, size=15)

    add_footer(slide, page_num, total, platform)
    return slide


def build_kpi_slide(prs: Presentation, eyebrow: str, title: str, kpis: list,
                    platform: str, page_num: int, total: int):
    """kpis = list of (big_value, label) tuples — up to 4."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_header(slide, title, eyebrow=eyebrow)

    n = len(kpis)
    card_w = Inches(2.7)
    gap = Inches(0.25)
    total_w = card_w * n + gap * (n - 1)
    start_left = (SLIDE_W - total_w) / 2

    top = Inches(2.6)
    card_h = Inches(2.4)

    for i, (big, label) in enumerate(kpis):
        left = start_left + (card_w + gap) * i
        # Card body
        card = add_rect(slide, left, top, card_w, card_h, TEXT_LIGHT)
        # Top color stripe
        add_rect(slide, left, top, card_w, Inches(0.12), BRAND_ACCENT)
        # Big value
        add_text(slide, left, top + Inches(0.4), card_w, Inches(1.2),
                 big, size=42, bold=True, color=BRAND_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_text(slide, left + Inches(0.2), top + Inches(1.55), card_w - Inches(0.4), Inches(0.8),
                 label, size=13, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    add_footer(slide, page_num, total, platform)
    return slide


def build_closing_slide(prs: Presentation, platform: str, message: str,
                        cta_lines: list, page_num: int, total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BRAND_DARK)

    add_rect(slide, Emu(0), Emu(0), Inches(0.6), SLIDE_H, BRAND_ACCENT)
    add_rect(slide, Inches(0.6), Emu(0), SLIDE_W - Inches(0.6), Inches(0.2), ACCENT_GOLD)

    add_text(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(1.2),
             "Thank you", size=54, bold=True, color=WHITE)
    add_rect(slide, Inches(1.2), Inches(3.4), Inches(2.0), Inches(0.08), ACCENT_GOLD)

    add_text(slide, Inches(1.2), Inches(3.7), Inches(11), Inches(0.8),
             message, size=20, color=TEXT_LIGHT)

    # CTA list
    add_bullets(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(2.0),
                cta_lines, size=18, color=TEXT_LIGHT, bullet_color=ACCENT_GOLD)

    add_footer(slide, page_num, total, platform)
    return slide


# ============================================================
# Deck definitions
# ============================================================

def define_internal_audit_deck():
    """Returns list of slide specs. First/last handled specially."""
    return {
        "platform": "Internal Audit Platform",
        "subtitle": "Risk-based, AI-augmented internal audit on a modern stack",
        "kpis": [
            ("50%", "Less time on workpaper review"),
            ("3 langs", "EN / AR (RTL) / LV"),
            ("Compliance-ready", "Aligned to your frameworks"),
            ("99.9%", "Production uptime SLA"),
        ],
        "kpi_eyebrow": "Outcomes at a glance",
        "kpi_title": "Why teams switch",
        "closing_msg": "Let's run a two-week pilot on a real engagement.",
        "closing_cta": [
            "Reach us at info@glimmora.ai",
            "30-minute discovery call",
            "60-minute live demo on your audit universe",
            "Two-week pilot on UAT with a real engagement",
            "Production go-live in 6 – 8 weeks",
        ],
        "sections": [
            # (section_label, section_title, [slide_specs])
            ("Section 1", "The Opportunity", [
                ("Eyebrow", "The Problem",
                 ["Internal audit is buried in spreadsheets, email, and tribal knowledge",
                  "Audit Committee asks 'Are we covering the right risks?' — answer takes weeks",
                  "CAPA closure is opaque; findings repeat year after year",
                  "External audit fees rise because the internal trail is hard to follow"]),
                ("Eyebrow", "Why Now",
                 ["Regulators demand demonstrable, risk-based assurance",
                  "Boards expect dashboard-grade reporting, not PDFs",
                  "AI has matured to safely augment workpapers, evidence review, and Q&A",
                  "Hybrid workforces need one source of truth — not a shared drive"]),
                ("Eyebrow", "Our Solution",
                 ["End-to-end Internal Audit lifecycle on one platform",
                  "Risk-based planning  →  fieldwork  →  findings  →  CAPA  →  reporting",
                  "AI-augmented at every stage — the human stays in control",
                  "Multi-tenant SaaS, or single-tenant in your cloud"]),
            ]),
            ("Section 2", "Product Tour", [
                ("Module Map", "The Audit Lifecycle in One View",
                 ["Audit Universe + Risk Universe feed the Annual Plan",
                  "Audit timelines are tracked within the Audit Universe itself — no parallel schedule to maintain",
                  "Engagements drive Fieldwork; Fieldwork generates Findings",
                  "Findings drive CAPA; CAPA closures feed Reports",
                  "Reports go straight to the Audit Committee — no manual assembly"]),
                ("Modules", "Eleven Modules, One Workspace",
                 ["Dashboard — execution, CAPA aging, risk-by-rating",
                  "Audit Universe + Risk Universe — your taxonomies, your scoring",
                  "Risk Identification — capture and triage emerging risks",
                  "Risk Register — single source of truth for risk exposure and treatment",
                  "Audit Planning — manual or AI-generated annual plan",
                  "Fieldwork — workpapers, evidence requests, findings",
                  "CAPA Tracking — owner, due date, AI review of closure",
                  "Reports + Document Library — board-ready outputs and searchable history"]),
                ("AI", "AI Capabilities — Embedded, Not Bolt-On",
                 ["AI-Recommended Audits — surface where to audit next",
                  "AI Workpapers — draft tests against process and risk",
                  "AI Evidence Review — verify completeness and relevance of uploaded evidence",
                  "Audit Q&A — chat with past engagements and evidence",
                  "AI CAPA Review — flags weak or off-topic closure evidence",
                  "Risk Suggestions — propose risks from process narratives",
                  "Translation — Arabic and Latvian for global teams"]),
            ]),
            ("Section 3", "Stakeholder Value", [
                ("Two-Column", "Who Wins and How",
                 None)  # handled separately below
            ]),
            ("Section 4", "Trust & Deployment", [
                ("Trust", "Security & Trust",
                 ["AES-256-GCM at rest, TLS 1.2+ in transit",
                  "RBAC with 8+ audit-specific roles enforced at three layers",
                  "Full activity log — every read, edit, approval",
                  "Aligns with frameworks that meet your organization's standards and requirements",
                  "Encryption kill switch + 90-day master-key rotation"]),
                ("Deploy", "Deployment Options",
                 ["Multi-tenant SaaS on DigitalOcean (default)",
                  "Single-tenant in customer cloud (AWS / Azure / GCP)",
                  "On-Premise — fully customer-hosted deployment supported",
                  "Hybrid — app in customer cloud, AI services on Glimmora backbone",
                  "UAT + Production environments included, gated by feature flags"]),
                ("Implementation", "Implementation Approach",
                 ["Week 1–2: Discovery — map audit universe and risk taxonomy",
                  "Week 3–4: Configuration — scoring, roles, processes, departments",
                  "Week 5–6: Pilot engagement with one audit team",
                  "Week 7–8: Migration of historical findings + CAPA",
                  "Week 9+: Full rollout with weekly office hours for 90 days"]),
            ]),
            ("Section 5", "Commercial & Roadmap", [
                ("Pricing", "Pricing Model",
                 ["Per-tenant annual subscription tied to user tiers",
                  "AI services metered separately for fair-use pricing",
                  "Single-tenant deployments quoted on infrastructure footprint",
                  "Implementation services billed per phase, fixed-fee",
                  "Cancellation within 15 days qualifies for a refund"]),
                ("Why Glimmora", "Why Glimmora",
                 ["Purpose-built for GRC — not a generic workflow tool",
                  "Founders bring decades of audit and consulting experience",
                  "Modern stack (Next.js 16, Python AI) — not a legacy product",
                  "Reference customers across BFSI, healthcare, and public sector"]),
                ("Roadmap", "Roadmap (12 Months)",
                 ["Continuous-control monitoring integrations",
                  "ServiceNow / Jira CAPA webhooks",
                  "Advanced analytics over multi-year audit history",
                  "Voice-driven evidence capture for field auditors"]),
            ]),
        ],
        # The two-column stakeholder slide content
        "stakeholder_two_col": {
            "left_heading": "Audit Function",
            "left_bullets": [
                "Audit Head — defensible methodology, board-ready reporting",
                "Audit Manager — workpaper review time cut in half",
                "Auditor — reusable processes, AI workpaper drafts, self-service evidence",
            ],
            "right_heading": "Business & Finance",
            "right_bullets": [
                "Auditee — one dashboard for every request, finding, and CAPA",
                "CFO — reduced external audit fees and TCO",
                "Board — live coverage view, not a quarterly PDF",
            ],
        },
    }


def define_grc_deck():
    return {
        "platform": "GRC Platform",
        "subtitle": "Governance, Risk, Compliance & Assets — one stack, one taxonomy",
        "kpis": [
            ("4 Pillars", "Org • Compliance • Risk • Assets"),
            ("Test once", "Satisfy many frameworks"),
            ("3 langs", "EN / AR (RTL) / LV"),
            ("ISO / SOC2", "GDPR / PCI mapped"),
        ],
        "kpi_eyebrow": "Outcomes at a glance",
        "kpi_title": "Why teams consolidate on Glimmora GRC",
        "closing_msg": "Let's run a two-week pilot on your frameworks and risks.",
        "closing_cta": [
            "30-minute discovery call",
            "60-minute live demo on your frameworks and risk register",
            "Two-week pilot on UAT with a real function",
            "Production go-live in 6 – 8 weeks",
        ],
        "sections": [
            ("Section 1", "The Opportunity", [
                ("Eyebrow", "The Problem",
                 ["Risk lives in spreadsheets, compliance in SharePoint, assets in CMDB",
                  "Same control is tested five times for five frameworks",
                  "BIA is a 200-page doc nobody reads until DR happens",
                  "Board asks 'What is our top risk?' and gets five different answers"]),
                ("Eyebrow", "Why Now",
                 ["Regulators require integrated risk + compliance + control evidence",
                  "Customers demand SOC 2 / ISO 27001 in their procurement RFPs",
                  "AI now usable for policy Q&A, risk suggestion, and translation",
                  "Hybrid workforces need a shared system of record"]),
                ("Eyebrow", "Our Solution",
                 ["One platform: Organization + Compliance + Risk + Assets",
                  "Shared taxonomy across all four pillars",
                  "AI-augmented and multilingual (EN / AR / LV)",
                  "Multi-tenant SaaS or single-tenant in customer cloud"]),
            ]),
            ("Section 2", "The Four Pillars", [
                ("Pillar 1", "Organization",
                 ["Organisation profile and context register",
                  "Process catalogue — the spine of the platform",
                  "Structured BIA per process with configurable methodology",
                  "BCP labels drive criticality across risk and audit"]),
                ("Pillar 2", "Compliance",
                 ["Framework library: ISO 27001, SOC 2, PCI, GDPR + custom",
                  "Controls mapped across frameworks — test once, satisfy many",
                  "Governance docs with review cadence",
                  "Evidence register, Exceptions, KPIs, Regulatory Intelligence",
                  "Statement of Applicability stays in sync automatically"]),
                ("Pillar 3", "Risk Management",
                 ["Risk register with assessment, response, review history",
                  "Risk-Control Matrix linking risks to mitigating controls",
                  "Configurable impact / probability / scoring ranges",
                  "AI suggests risks from process and governance narratives"]),
                ("Pillar 4", "Asset Management",
                 ["Asset Inventory with custom attributes and classification",
                  "My Inventory view for custodians",
                  "Asset ↔ Risk ↔ Control linking for full traceability",
                  "Reports for management and audit review"]),
            ]),
            ("Section 3", "Differentiators", [
                ("AI", "AI Capabilities",
                 ["RAG over governance documents — chat with your policies",
                  "Risk suggestions from process narratives",
                  "Document ingestion with OCR + embeddings",
                  "Dynamic translation of user-entered data on create/edit"]),
                ("Two-Column", "Stakeholder Value",
                 None),  # handled separately
            ]),
            ("Section 4", "Trust & Deployment", [
                ("Trust", "Security & Trust",
                 ["AES-256-GCM at rest, TLS 1.2+ in transit",
                  "RBAC with department-scoped roles",
                  "Activity log for every change",
                  "Maps to ISO 27001, SOC 2, GDPR, PCI",
                  "Encryption kill switch + 90-day master-key rotation"]),
                ("Deploy", "Deployment Options",
                 ["Multi-tenant SaaS on DigitalOcean (default)",
                  "Single-tenant in customer cloud (AWS / Azure / GCP)",
                  "Hybrid — app in customer cloud, AI services on Glimmora backbone",
                  "UAT + Production environments included with feature-flag gating"]),
                ("Implementation", "Implementation Approach",
                 ["Week 1–2: Discovery and process / framework scoping",
                  "Week 3–4: Configuration — frameworks, taxonomy, roles, departments",
                  "Week 5–6: Pilot with one function (e.g., IT or Finance)",
                  "Week 7–8: Migration of historical risks, controls, evidence",
                  "Week 9+: Full rollout + 90 days of weekly office hours"]),
            ]),
            ("Section 5", "Commercial & Roadmap", [
                ("Pricing", "Pricing Model",
                 ["Per-tenant annual subscription tied to user tiers",
                  "AI services metered separately for fair-use pricing",
                  "Single-tenant deployments quoted on infrastructure footprint",
                  "Implementation services billed per phase, fixed-fee"]),
                ("Why Glimmora", "Why Glimmora",
                 ["GRC purpose-built — not a generic workflow tool",
                  "Founders bring decades of GRC and consulting experience",
                  "Modern stack (Next.js 16, Python AI)",
                  "Native multilingual including Arabic RTL"]),
                ("Roadmap", "Roadmap (12 Months)",
                 ["Continuous-control monitoring connectors",
                  "Native mobile evidence capture",
                  "ITSM webhooks (ServiceNow / Jira)",
                  "Quantitative risk (Monte Carlo) module",
                  "Vendor-risk deep integration with TPRM platform"]),
            ]),
        ],
        "stakeholder_two_col": {
            "left_heading": "Risk & Compliance Leaders",
            "left_bullets": [
                "CRO — live risk dashboard, shared taxonomy",
                "CCO — multi-framework mapping, evidence reuse",
                "BCM Lead — structured BIA tied to process catalogue",
            ],
            "right_heading": "Business & Finance",
            "right_bullets": [
                "Control / Process Owner — clear ownership, less duplicate work",
                "CISO — encryption, RBAC, secure SDLC",
                "CFO — fewer findings, lower TCO than separate point tools",
            ],
        },
    }


# ---------- Build a deck ----------

def build_deck(spec: dict, out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    platform = spec["platform"]

    # First pass — count total slides for the footer
    # Structure: cover + KPI + (section divider + content slides) per section + closing
    total = 1  # cover
    total += 1  # KPI
    for _, _, items in spec["sections"]:
        total += 1  # section divider
        total += len(items)
    total += 1  # closing

    page = 1

    # Cover
    build_cover_slide(prs, platform, spec["subtitle"])
    page += 1  # cover doesn't show footer/page num, but we count it

    # KPI snapshot right after cover
    build_kpi_slide(prs, spec["kpi_eyebrow"], spec["kpi_title"], spec["kpis"],
                    platform, page, total)
    page += 1

    # Sections
    for section_label, section_title, items in spec["sections"]:
        build_section_divider(prs, section_label, section_title, platform, page, total)
        page += 1

        for spec_tuple in items:
            eyebrow, title, bullets = spec_tuple

            if bullets is None and title in ("Who Wins and How", "Stakeholder Value"):
                tc = spec["stakeholder_two_col"]
                build_two_col_slide(
                    prs, eyebrow, title,
                    tc["left_heading"], tc["left_bullets"],
                    tc["right_heading"], tc["right_bullets"],
                    platform, page, total,
                )
            else:
                build_content_slide(prs, eyebrow, title, bullets,
                                    platform, page, total)
            page += 1

    # Closing
    build_closing_slide(prs, platform, spec["closing_msg"], spec["closing_cta"],
                        page, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return total


def main():
    root = Path(__file__).resolve().parent.parent
    out_dir = root / "docs" / "platform-docs"

    print(f"Output directory: {out_dir}")

    jobs = [
        (define_internal_audit_deck(), out_dir / "InternalAudit-Presentation.pptx"),
        (define_grc_deck(),            out_dir / "GRC-Presentation.pptx"),
    ]

    for spec, path in jobs:
        n = build_deck(spec, path)
        print(f"  wrote {path.name}  ({n} slides, {path.stat().st_size // 1024} KB)")

    # Remove the older .docx presentations so the deck-format files don't conflict
    for legacy in [
        out_dir / "InternalAudit-Presentation.docx",
        out_dir / "GRC-Presentation.docx",
    ]:
        if legacy.exists():
            legacy.unlink()
            print(f"  removed legacy {legacy.name}")

    print(f"\nDone. {len(jobs)} presentations generated in {out_dir}")


if __name__ == "__main__":
    main()
